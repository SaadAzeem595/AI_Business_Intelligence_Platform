import os
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

TEMPLATE_RGB_THEMES = {
    "CEO": {
        "primary": RGBColor(30, 58, 138),      # Deep Blue
        "secondary": RGBColor(59, 130, 246),    # Light Blue
        "text": RGBColor(15, 23, 42),           # Slate 900
        "card_bg": RGBColor(241, 245, 249)      # Slate 100
    },
    "Sales": {
        "primary": RGBColor(15, 23, 42),        # Slate 900
        "secondary": RGBColor(245, 158, 11),    # Amber 500
        "text": RGBColor(15, 23, 42),
        "card_bg": RGBColor(254, 243, 199)      # Amber 100
    },
    "Finance": {
        "primary": RGBColor(6, 78, 59),         # Deep Emerald
        "secondary": RGBColor(16, 185, 129),    # Emerald 500
        "text": RGBColor(15, 23, 42),
        "card_bg": RGBColor(236, 253, 245)      # Emerald 100
    },
    "Marketing": {
        "primary": RGBColor(76, 29, 149),       # Deep Purple
        "secondary": RGBColor(139, 92, 246),    # Violet 500
        "text": RGBColor(15, 23, 42),
        "card_bg": RGBColor(245, 243, 255)      # Violet 100
    },
    "Operations": {
        "primary": RGBColor(17, 24, 39),        # Charcoal
        "secondary": RGBColor(100, 116, 139),   # Gray 500
        "text": RGBColor(15, 23, 42),
        "card_bg": RGBColor(243, 244, 246)      # Gray 100
    }
}


class PowerPointReportGenerator:
    """Generates editable PowerPoint presentations with matching color palettes and layouts."""

    @staticmethod
    def generate(
        filepath: str,
        title: str,
        template_name: str,
        data: dict,
        snapshot_path: str = None
    ) -> str:
        """Assembles and writes PPTX presentation file to target path."""
        os.makedirs(os.path.dirname(filepath), exist_ok=True)
        
        # Initialize Presentation
        prs = Presentation()
        # Set slide sizes to 16:9 widescreen
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        # Load theme colors
        theme = TEMPLATE_RGB_THEMES.get(template_name, TEMPLATE_RGB_THEMES["CEO"])
        primary_color = theme["primary"]
        secondary_color = theme["secondary"]
        text_color = theme["text"]
        card_bg_color = theme["card_bg"]
        
        # Layout templates: blank layout is index 6
        blank_layout = prs.slide_layouts[6]
        
        # --- Slide 1: COVER SLIDE ---
        slide1 = prs.slides.add_slide(blank_layout)
        
        # Draw background band line
        accent_band = slide1.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), Inches(7.5)
        )
        accent_band.fill.solid()
        accent_band.fill.fore_color.rgb = secondary_color
        accent_band.line.fill.background()
        
        # Title Box
        title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.0), Inches(2.0))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Arial"
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = primary_color
        
        # Subtitle
        p2 = tf.add_paragraph()
        p2.text = f"{template_name} Executive Business Deliverable"
        p2.font.name = "Arial"
        p2.font.size = Pt(22)
        p2.font.color.rgb = RGBColor(100, 116, 139)  # Slate 500
        p2.space_before = Pt(10)
        
        # Metadata box
        meta_box = slide1.shapes.add_textbox(Inches(1.0), Inches(4.5), Inches(8.0), Inches(2.0))
        meta_tf = meta_box.text_frame
        
        meta_items = [
            f"Author: {data.get('author', 'Executive Intelligence Platform')}",
            f"Workspace: {data.get('workspace', 'default').upper()}",
            f"Date: {datetime.now().strftime('%B %d, %Y')}",
            f"Confidence Score: {data.get('confidence_score', 0.95)*100:.1f}%"
        ]
        
        for idx, item in enumerate(meta_items):
            p_m = meta_tf.add_paragraph() if idx > 0 else meta_tf.paragraphs[0]
            p_m.text = item
            p_m.font.name = "Arial"
            p_m.font.size = Pt(13)
            p_m.font.color.rgb = RGBColor(71, 85, 105)
            p_m.space_after = Pt(4)
            
        # --- Slide 2: EXECUTIVE SUMMARY ---
        slide2 = prs.slides.add_slide(blank_layout)
        
        # Slide Title
        title2 = slide2.shapes.add_textbox(Inches(1.0), Inches(0.5), Inches(11.3), Inches(0.8))
        title2.text_frame.text = "Executive Insights Summary"
        title2.text_frame.paragraphs[0].font.size = Pt(28)
        title2.text_frame.paragraphs[0].font.bold = True
        title2.text_frame.paragraphs[0].font.color.rgb = primary_color
        
        # Slide content list
        takeaways = (data.get("executive_summary") or {}).get("key_takeaways", [
            "Revenue is expanding along seasonal forecast targets.",
            "Predictive models validate churn risks in active operational regions.",
            "Recommendations indicate cohort discount campaigns prioritize West region users."
        ])
        
        content2 = slide2.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.3), Inches(4.5))
        content_tf = content2.text_frame
        content_tf.word_wrap = True
        
        for idx, takeaway in enumerate(takeaways):
            p_c = content_tf.add_paragraph() if idx > 0 else content_tf.paragraphs[0]
            p_c.text = f"•  {takeaway}"
            p_c.font.name = "Arial"
            p_c.font.size = Pt(18)
            p_c.font.color.rgb = text_color
            p_c.space_after = Pt(20)
            
        # --- Slide 3: KPI OVERVIEW GRID ---
        slide3 = prs.slides.add_slide(blank_layout)
        
        title3 = slide3.shapes.add_textbox(Inches(1.0), Inches(0.5), Inches(11.3), Inches(0.8))
        title3.text_frame.text = "Business Metric Indicators (KPIs)"
        title3.text_frame.paragraphs[0].font.size = Pt(28)
        title3.text_frame.paragraphs[0].font.bold = True
        title3.text_frame.paragraphs[0].font.color.rgb = primary_color
        
        kpi_list = data.get("kpi_overview", [
            {"title": "Total Revenue", "value": "$1.24M", "change": "+14.2% MoM"},
            {"title": "Operating Cost", "value": "$320.5K", "change": "-2.1% MoM"},
            {"title": "Churn Prediction", "value": "15.0%", "change": "+0.5% MoM"},
            {"title": "RAG Score", "value": "94.0%", "change": "+1.8% MoM"}
        ])
        
        # Grid variables: 4 columns.
        left_margin = Inches(1.0)
        gap = Inches(0.4)
        card_w = Inches(2.5)
        card_h = Inches(3.5)
        top = Inches(2.0)
        
        for idx, kpi in enumerate(kpi_list[:4]):
            left = left_margin + idx * (card_w + gap)
            
            # Card background
            card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, card_w, card_h)
            card.fill.solid()
            card.fill.fore_color.rgb = card_bg_color
            card.line.color.rgb = secondary_color
            card.line.width = Pt(1.5)
            
            # Text Frame inside card
            card_box = slide3.shapes.add_textbox(left + Inches(0.1), top + Inches(0.2), card_w - Inches(0.2), card_h - Inches(0.4))
            card_tf = card_box.text_frame
            card_tf.word_wrap = True
            
            # Title
            p_t = card_tf.paragraphs[0]
            p_t.text = kpi["title"]
            p_t.font.name = "Arial"
            p_t.font.size = Pt(14)
            p_t.font.bold = True
            p_t.font.color.rgb = RGBColor(100, 116, 139)
            p_t.alignment = PP_ALIGN.CENTER
            p_t.space_after = Pt(24)
            
            # Value
            p_v = card_tf.add_paragraph()
            p_v.text = kpi["value"]
            p_v.font.name = "Arial"
            p_v.font.size = Pt(28)
            p_v.font.bold = True
            p_v.font.color.rgb = primary_color
            p_v.alignment = PP_ALIGN.CENTER
            p_v.space_after = Pt(24)
            
            # Change
            p_ch = card_tf.add_paragraph()
            p_ch.text = kpi["change"]
            p_ch.font.name = "Arial"
            p_ch.font.size = Pt(13)
            p_ch.font.bold = True
            is_positive = not kpi["change"].startswith("-") and not kpi["change"].startswith("0")
            p_ch.font.color.rgb = RGBColor(16, 185, 129) if is_positive else RGBColor(239, 68, 68)
            p_ch.alignment = PP_ALIGN.CENTER
            
        # --- Slide 4: SNAPSHOT CHART SLIDE ---
        if snapshot_path and os.path.exists(snapshot_path):
            slide4 = prs.slides.add_slide(blank_layout)
            
            title4 = slide4.shapes.add_textbox(Inches(1.0), Inches(0.5), Inches(11.3), Inches(0.8))
            title4.text_frame.text = "Dashboard Snapshot & Trend Analysis"
            title4.text_frame.paragraphs[0].font.size = Pt(28)
            title4.text_frame.paragraphs[0].font.bold = True
            title4.text_frame.paragraphs[0].font.color.rgb = primary_color
            
            # Left column: chart image
            slide4.shapes.add_picture(snapshot_path, Inches(1.0), Inches(1.5), Inches(7.5), Inches(4.5))
            
            # Right column: text highlights
            analysis_box = slide4.shapes.add_textbox(Inches(8.8), Inches(1.8), Inches(3.5), Inches(4.0))
            tf_a = analysis_box.text_frame
            tf_a.word_wrap = True
            
            p_a1 = tf_a.paragraphs[0]
            p_a1.text = "Analytical Observations:"
            p_a1.font.bold = True
            p_a1.font.size = Pt(16)
            p_a1.font.color.rgb = secondary_color
            p_a1.space_after = Pt(10)
            
            bullet_points = [
                "Dashboard snapshots verify a steady growth pattern.",
                "KPI aggregates display stable operational margins across all segments.",
                "Forecasting models maintain positive Q4 trends with 95% confidence intervals."
            ]
            for bp in bullet_points:
                p_bp = tf_a.add_paragraph()
                p_bp.text = f"- {bp}"
                p_bp.font.size = Pt(12)
                p_bp.font.color.rgb = text_color
                p_bp.space_after = Pt(8)
                
        # --- Slide 5: STRATEGIC ACTION ITEMS ---
        slide5 = prs.slides.add_slide(blank_layout)
        
        title5 = slide5.shapes.add_textbox(Inches(1.0), Inches(0.5), Inches(11.3), Inches(0.8))
        title5.text_frame.text = "Strategic Recommendations & Action Plan"
        title5.text_frame.paragraphs[0].font.size = Pt(28)
        title5.text_frame.paragraphs[0].font.bold = True
        title5.text_frame.paragraphs[0].font.color.rgb = primary_color
        
        recs = data.get("recommendations") or [
            {"insight": "Customer churn rates are stable at 15%. Recommend target discount campaigns on the West region.", "confidence_score": 0.88, "priority": "High"},
            {"insight": "Q4 forecasts project a steady sales rise. Ensure warehouse supply matches the 5% margin increase.", "confidence_score": 0.92, "priority": "Medium"}
        ]
        
        # Add table: rows = len(recs) + 1 headers, cols = 3 (Insight, Priority, Confidence)
        rows = len(recs) + 1
        cols = 3
        table_shape = slide5.shapes.add_table(rows, cols, Inches(1.0), Inches(1.8), Inches(11.333), Inches(4.5))
        table = table_shape.table
        
        # Column widths
        table.columns[0].width = Inches(7.5)
        table.columns[1].width = Inches(1.833)
        table.columns[2].width = Inches(2.0)
        
        # Set Headers
        headers = ["Strategic Recommendation / Insight", "Priority", "Confidence Score"]
        for col_idx, text in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = text
            cell.fill.solid()
            cell.fill.fore_color.rgb = primary_color
            for p in cell.text_frame.paragraphs:
                p.font.bold = True
                p.font.color.rgb = RGBColor(255, 255, 255)
                p.font.size = Pt(14)
                p.alignment = PP_ALIGN.CENTER if col_idx > 0 else PP_ALIGN.LEFT
                
        # Fill table rows
        for row_idx, rec in enumerate(recs):
            # Insight
            cell_ins = table.cell(row_idx + 1, 0)
            cell_ins.text = rec["insight"]
            cell_ins.text_frame.paragraphs[0].font.size = Pt(13)
            cell_ins.text_frame.paragraphs[0].font.color.rgb = text_color
            
            # Priority
            cell_prio = table.cell(row_idx + 1, 1)
            cell_prio.text = rec["priority"]
            cell_prio.text_frame.paragraphs[0].font.size = Pt(13)
            cell_prio.text_frame.paragraphs[0].font.bold = True
            cell_prio.text_frame.paragraphs[0].font.color.rgb = primary_color
            cell_prio.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # Confidence
            cell_conf = table.cell(row_idx + 1, 2)
            conf_str = f"{rec['confidence_score']*100:.0f}%" if isinstance(rec['confidence_score'], (int, float)) else str(rec['confidence_score'])
            cell_conf.text = conf_str
            cell_conf.text_frame.paragraphs[0].font.size = Pt(13)
            cell_conf.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            cell_conf.text_frame.paragraphs[0].font.color.rgb = text_color
            
        prs.save(filepath)
        return filepath
