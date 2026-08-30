import io
import json
import logging
from abc import ABC, abstractmethod
from typing import Dict, Type, Optional
import pandas as pd
from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation as PptxPresentation
from bs4 import BeautifulSoup

from app.features.rag.ingestion.ocr import BaseOCRProvider, MockOCRProvider

logger = logging.getLogger(__name__)

class BaseParser(ABC):
    @abstractmethod
    def parse(self, file_bytes: bytes, filename: str, ocr_provider: Optional[BaseOCRProvider] = None) -> str:
        """Parses document bytes and extracts raw text."""
        pass


class TextParser(BaseParser):
    def parse(self, file_bytes: bytes, filename: str, ocr_provider: Optional[BaseOCRProvider] = None) -> str:
        return file_bytes.decode("utf-8", errors="ignore")


class MarkdownParser(BaseParser):
    def parse(self, file_bytes: bytes, filename: str, ocr_provider: Optional[BaseOCRProvider] = None) -> str:
        return file_bytes.decode("utf-8", errors="ignore")


class PDFParser(BaseParser):
    def parse(self, file_bytes: bytes, filename: str, ocr_provider: Optional[BaseOCRProvider] = None) -> str:
        pdf_file = io.BytesIO(file_bytes)
        reader = PdfReader(pdf_file)
        text_parts = []
        
        for i, page in enumerate(reader.pages):
            page_text = page.extract_text()
            if page_text:
                text_parts.append(page_text)
                
        full_text = "\n".join(text_parts).strip()
        
        # If the extracted text is very short/empty, run OCR (suggesting scanned PDF)
        if len(full_text) < 50:
            logger.info(f"PDF '{filename}' appears scanned. Running pluggable OCR...")
            full_text = ocr_provider.extract_text(file_bytes)
            
        return full_text


class DocxParser(BaseParser):
    def parse(self, file_bytes: bytes, filename: str, ocr_provider: Optional[BaseOCRProvider] = None) -> str:
        doc = DocxDocument(io.BytesIO(file_bytes))
        text_parts = []
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_parts.append(paragraph.text)
        for table in doc.tables:
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_text:
                    text_parts.append(" | ".join(row_text))
        return "\n".join(text_parts)


class PptxParser(BaseParser):
    def parse(self, file_bytes: bytes, filename: str, ocr_provider: Optional[BaseOCRProvider] = None) -> str:
        prs = PptxPresentation(io.BytesIO(file_bytes))
        text_parts = []
        for i, slide in enumerate(prs.slides):
            text_parts.append(f"--- Slide {i+1} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text)
        return "\n".join(text_parts)


class HtmlParser(BaseParser):
    def parse(self, file_bytes: bytes, filename: str, ocr_provider: Optional[BaseOCRProvider] = None) -> str:
        soup = BeautifulSoup(file_bytes, "html.parser")
        # Remove script and style elements
        for script in soup(["script", "style"]):
            script.decompose()
        return soup.get_text(separator="\n")


class CsvParser(BaseParser):
    @staticmethod
    def format_dataframe(df: pd.DataFrame, filename: str, sheet_name: str = None) -> str:
        if df.empty:
            return ""
            
        cols = [str(c).strip() for c in df.columns]
        cols_str = " | ".join(cols)
        sheet_info = f" (Sheet: {sheet_name})" if sheet_name else ""
        
        header_block = f"[TABULAR_DATA: {filename}{sheet_info}]\n[SCHEMA: {cols_str}]"
        
        records = df.to_dict(orient="records")
        rows_formatted = []
        for idx, rec in enumerate(records):
            row_kvs = [f"{col}: {rec[col]}" for col in cols if col in rec and pd.notna(rec[col])]
            row_str = f"Row {idx + 1} -> " + " | ".join(row_kvs)
            rows_formatted.append(row_str)
            
        return header_block + "\n" + "\n".join(rows_formatted)

    def parse(self, file_bytes: bytes, filename: str, ocr_provider: Optional[BaseOCRProvider] = None) -> str:
        df = pd.read_csv(io.BytesIO(file_bytes))
        return self.format_dataframe(df, filename)


class ExcelParser(BaseParser):
    def parse(self, file_bytes: bytes, filename: str, ocr_provider: Optional[BaseOCRProvider] = None) -> str:
        # Load excel file (handles xls and xlsx)
        xls = pd.ExcelFile(io.BytesIO(file_bytes))
        text_parts = []
        for sheet_name in xls.sheet_names:
            df = pd.read_excel(xls, sheet_name=sheet_name)
            formatted = CsvParser.format_dataframe(df, filename, sheet_name=sheet_name)
            if formatted:
                text_parts.append(formatted)
        return "\n\n".join(text_parts)


class JsonParser(BaseParser):
    def parse(self, file_bytes: bytes, filename: str, ocr_provider: Optional[BaseOCRProvider] = None) -> str:
        data = json.loads(file_bytes.decode("utf-8", errors="ignore"))
        return json.dumps(data, indent=2)


class DocumentParserService:
    """Registry-based service that resolves and executes correct file parser."""
    
    _parsers: Dict[str, Type[BaseParser]] = {
        "txt": TextParser,
        "text": TextParser,
        "md": MarkdownParser,
        "markdown": MarkdownParser,
        "pdf": PDFParser,
        "docx": DocxParser,
        "pptx": PptxParser,
        "html": HtmlParser,
        "htm": HtmlParser,
        "csv": CsvParser,
        "xlsx": ExcelParser,
        "xls": ExcelParser,
        "json": JsonParser
    }
    
    def __init__(self, ocr_provider: BaseOCRProvider = None):
        self.ocr_provider = ocr_provider or MockOCRProvider()
        
    def parse_file(self, file_bytes: bytes, filename: str) -> str:
        ext = filename.split(".")[-1].lower() if "." in filename else "txt"
        parser_cls = self._parsers.get(ext, TextParser)
        
        logger.info(f"Parsing file '{filename}' using {parser_cls.__name__}")
        parser = parser_cls()
        return parser.parse(file_bytes, filename, self.ocr_provider)
